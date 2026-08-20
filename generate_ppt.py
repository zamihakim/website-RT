"""
Generate Presentasi Aplikasi Iuran RT dengan UI/UX Wireframe Mockups
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
TEAL       = RGBColor(0x0F, 0x76, 0x6E)
DARK_TEAL  = RGBColor(0x06, 0x5F, 0x46)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1F, 0x29, 0x37)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
LIGHTER    = RGBColor(0xE2, 0xE8, 0xF0)
BLUE       = RGBColor(0x3B, 0x82, 0xF6)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
YELLOW_BG  = RGBColor(0xFF, 0xC1, 0x07)
RED        = RGBColor(0xEF, 0x44, 0x44)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
DARK_BG    = RGBColor(0x1E, 0x29, 0x3B)
SIDEBAR_BG = RGBColor(0x34, 0x3A, 0x40)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def set_shape_text(shape, text, font_size=12, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = 'Segoe UI'
    tf.paragraphs[0].alignment = alignment

def add_slide_title(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.3), TEAL)
    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7),
                title, font_size=32, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.4),
                    subtitle, font_size=16, color=RGBColor(0xCC, 0xFB, 0xF1))

def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 13

# ================================================================
# SLIDE 1 — Cover
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_TEAL)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), TEAL)
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            "Aplikasi Iuran RT", font_size=52, bold=True, color=WHITE)
add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.7),
            "Sistem Pencatatan Pemasukan & Pengeluaran Kas RT",
            font_size=24, color=RGBColor(0xA7, 0xF3, 0xD0))
add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
            "Mata Kuliah: Pemrograman Web Lanjut (PWL)",
            font_size=18, color=RGBColor(0xCC, 0xFB, 0xF1))
add_rect(slide, Inches(1.5), Inches(5.0), Inches(3), Inches(0.04), WHITE)
add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.5),
            "Nama  —  NIM  —  Kelas", font_size=16, color=WHITE)

# ================================================================
# SLIDE 2 — Latar Belakang & Tujuan
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Latar Belakang & Tujuan")

# Kolom kiri — Masalah
add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.5),
            "Masalah", font_size=22, bold=True, color=RED)
problems = [
    "• Pencatatan iuran & pengeluaran masih manual (buku)",
    "• Warga tidak tahu status pembayarannya",
    "• Rekap & laporan bulanan lambat dan rawan salah",
    "• Sulit memantau warga yang menunggak",
]
for i, p in enumerate(problems):
    add_textbox(slide, Inches(1.2), Inches(2.7 + i*0.65), Inches(5), Inches(0.6),
                p, font_size=16, color=BLACK)

# Kolom kanan — Tujuan
add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.6), Inches(5.0),
                 RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_textbox(slide, Inches(7.2), Inches(1.9), Inches(4), Inches(0.5),
            "Tujuan", font_size=22, bold=True, color=DARK_TEAL)
goals = [
    "• Mencatat iuran bulanan warga secara digital",
    "• Memantau siapa sudah/belum/macet membayar",
    "• Mencatat pengeluaran (Kas, Sosial, Konsumsi)",
    "• Menghasilkan laporan bulanan otomatis",
]
for i, g in enumerate(goals):
    add_textbox(slide, Inches(7.2), Inches(2.7 + i*0.65), Inches(5), Inches(0.6),
                g, font_size=16, color=BLACK)
add_page_number(slide, 2, TOTAL_SLIDES)

# ================================================================
# SLIDE 3 — Deskripsi Kebutuhan
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Deskripsi Kebutuhan Sistem")

# Table peran
table_data = [
    ["Peran", "Hak Akses"],
    ["Pengurus", "Semua menu pengelolaan + laporan"],
    ["Warga", "Tagihan, bayar iuran, riwayat pembayaran"],
]
tbl = slide.shapes.add_table(3, 2, Inches(0.8), Inches(1.7), Inches(6), Inches(2)).table
tbl.columns[0].width = Inches(2)
tbl.columns[1].width = Inches(4)
for r, row in enumerate(table_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = 'Segoe UI'
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL if r == 0 else (LIGHT_GRAY if r % 2 == 1 else WHITE)

# Alur pemasukan / pengeluaran
add_rounded_rect(slide, Inches(7.2), Inches(1.7), Inches(5.4), Inches(2.5),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(7.5), Inches(1.9), Inches(4), Inches(0.4),
            "Alur Keuangan", font_size=18, bold=True, color=TEAL)
add_textbox(slide, Inches(7.5), Inches(2.5), Inches(5), Inches(0.4),
            "Pemasukan = Iuran warga", font_size=15, color=GREEN)
add_textbox(slide, Inches(7.5), Inches(3.0), Inches(5), Inches(0.4),
            "Pengeluaran = Kas | Sosial | Konsumsi", font_size=15, color=RED)
add_textbox(slide, Inches(7.5), Inches(3.5), Inches(5), Inches(0.4),
            "Saldo = Total Pemasukan − Total Pengeluaran", font_size=15, color=BLUE)

# Fitur utama
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.5),
            "Fitur Utama", font_size=20, bold=True, color=TEAL)
features = [
    "✅  Dashboard ringkasan kas RT",
    "✅  Monitoring pembayaran warga",
    "✅  Tagihan iuran warga",
    "✅  Riwayat pembayaran",
    "✅  Laporan kas bulanan",
    "✅  Pengeluaran RT per kategori",
]
for i, f in enumerate(features):
    col = i // 3
    row = i % 3
    add_textbox(slide, Inches(0.8 + col*3.2), Inches(4.9 + row*0.55), Inches(3.2), Inches(0.5),
                f, font_size=14, color=BLACK)
add_page_number(slide, 3, TOTAL_SLIDES)

# ================================================================
# SLIDE 4 — Task 1: Pemilihan Template
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Task 1 — Pemilihan Template", "AdminLTE 3 (v3.2.0) • Bootstrap 4 • Lisensi MIT")

# Alasan
add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.0),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.5),
            "Alasan Pemilihan", font_size=20, bold=True, color=TEAL)
reasons = [
    "• Struktur layout jelas: navbar, sidebar, content, footer",
    "• Cocok untuk aplikasi admin/dashboard",
    "• Komponen siap pakai (info box, tabel, form, chart)",
    "• Responsif & banyak referensi tutorial",
    "• Open source, lisensi MIT",
]
for i, r in enumerate(reasons):
    add_textbox(slide, Inches(1.2), Inches(2.7 + i*0.6), Inches(5.2), Inches(0.5),
                r, font_size=15, color=BLACK)

# Struktur folder
add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.0),
                 DARK_BG, None)
add_textbox(slide, Inches(7.4), Inches(1.9), Inches(4), Inches(0.4),
            "Struktur Template", font_size=18, bold=True, color=GREEN)
folder_text = (
    "template_adminlte/\n"
    "├── dist/           (css, js, img)\n"
    "├── plugins/        (bootstrap, jquery)\n"
    "├── pages/\n"
    "├── starter.html\n"
    "└── index.html"
)
add_textbox(slide, Inches(7.4), Inches(2.5), Inches(5), Inches(3),
            folder_text, font_size=14, color=RGBColor(0xE2, 0xE8, 0xF0),
            font_name='Consolas')

# Alternatif
add_textbox(slide, Inches(7.4), Inches(5.0), Inches(5), Inches(0.4),
            "Alternatif yang dipertimbangkan:", font_size=13, bold=True, color=GRAY)
add_textbox(slide, Inches(7.4), Inches(5.4), Inches(5), Inches(1),
            "Stisla • SB Admin 2 • Argon Dashboard", font_size=13, color=GRAY)
add_page_number(slide, 4, TOTAL_SLIDES)

# ================================================================
# SLIDE 5 — Task 2: Slicing Layout
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Task 2 — Slicing Layout", "CodeIgniter 4 • extend() / section() / include()")

# Folder structure
add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.2),
                 DARK_BG, None)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.4),
            "Struktur Views", font_size=18, bold=True, color=GREEN)
views_text = (
    "app/Views/\n"
    "├── layout/\n"
    "│   └── layout.php      ← File Layout\n"
    "├── components/\n"
    "│   ├── header.php      ← <head>, CSS\n"
    "│   ├── navbar.php      ← top bar\n"
    "│   ├── sidebar.php     ← menu navigasi\n"
    "│   └── footer.php      ← footer + script\n"
    "└── pages/\n"
    "    ├── dashboard.php\n"
    "    ├── pembayaran.php\n"
    "    ├── laporan.php\n"
    "    ├── warga_tagihan.php\n"
    "    └── warga_history.php"
)
add_textbox(slide, Inches(1.2), Inches(2.5), Inches(5), Inches(4),
            views_text, font_size=13, color=RGBColor(0xE2, 0xE8, 0xF0),
            font_name='Consolas')

# Konsep
add_rounded_rect(slide, Inches(6.7), Inches(1.7), Inches(5.8), Inches(2.3),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(7.0), Inches(1.9), Inches(4), Inches(0.4),
            "Konsep Layout (layout.php)", font_size=16, bold=True, color=TEAL)
layout_code = (
    "<?= $this->include('components/header') ?>\n"
    "<?= $this->include('components/navbar') ?>\n"
    "<?= $this->include('components/sidebar') ?>\n"
    "<div class=\"content-wrapper\">\n"
    "  <?= $this->renderSection('content') ?>\n"
    "</div>\n"
    "<?= $this->include('components/footer') ?>"
)
add_textbox(slide, Inches(7.0), Inches(2.4), Inches(5.3), Inches(1.5),
            layout_code, font_size=11, color=RGBColor(0x33, 0x41, 0x55),
            font_name='Consolas')

# Contoh halaman
add_rounded_rect(slide, Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.6),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(7.0), Inches(4.5), Inches(4), Inches(0.4),
            "Contoh Page (dashboard.php)", font_size=16, bold=True, color=TEAL)
page_code = (
    "<?= $this->extend('layout/layout') ?>\n\n"
    "<?= $this->section('content') ?>\n"
    "  <!-- info box & tabel -->\n"
    "<?= $this->endSection() ?>"
)
add_textbox(slide, Inches(7.0), Inches(5.0), Inches(5.3), Inches(1.5),
            page_code, font_size=11, color=RGBColor(0x33, 0x41, 0x55),
            font_name='Consolas')
add_page_number(slide, 5, TOTAL_SLIDES)

# ================================================================
# SLIDE 6 — UI/UX: Dashboard Pengurus (WIREFRAME)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "UI/UX — Dashboard Pengurus", "Ringkasan kas, rekap pembayaran, pengeluaran bulanan")

# ── Sidebar ──
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(5.6), SIDEBAR_BG)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.8), DARK_TEAL)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.7),
            "🏠 Iuran RT", font_size=14, bold=True, color=WHITE)

sidebar_items = [
    ("📊 Dashboard", True),
    ("👥 Kelola Warga", False),
    ("⚙️ Pengaturan Iuran", False),
    ("💰 Monitoring Bayar", False),
    ("⚠️ Pembayaran Macet", False),
    ("📄 Pengeluaran RT", False),
    ("📈 Laporan Bulanan", False),
]
for i, (item, active) in enumerate(sidebar_items):
    y = Inches(2.4 + i * 0.55)
    bg = TEAL if active else SIDEBAR_BG
    add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.48), bg)
    add_textbox(slide, Inches(0.7), y, Inches(2), Inches(0.45),
                item, font_size=11, bold=active, color=WHITE)

# ── Top navbar ──
add_rect(slide, Inches(2.8), Inches(1.5), Inches(10), Inches(0.7), WHITE, LIGHTER, Pt(1))
add_textbox(slide, Inches(3.0), Inches(1.55), Inches(3), Inches(0.5),
            "Dashboard", font_size=16, bold=True, color=BLACK)
add_textbox(slide, Inches(11), Inches(1.55), Inches(1.5), Inches(0.5),
            "👤 Pengurus", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

# ── Info Boxes Row ──
box_data = [
    ("Rp 12.500.000", "Total Kas RT", BLUE, "💰"),
    ("38/50", "Sudah Bayar", GREEN, "✅"),
    ("12", "Belum Bayar", YELLOW_BG, "⏳"),
    ("Rp 3.250.000", "Pengeluaran", RED, "📄"),
]
for i, (val, label, color, icon) in enumerate(box_data):
    x = Inches(2.8 + i * 2.55)
    add_rounded_rect(slide, x, Inches(2.4), Inches(2.3), Inches(1.3), WHITE, LIGHTER)
    add_rect(slide, x, Inches(2.4), Inches(2.3), Inches(0.08), color)
    add_textbox(slide, x + Inches(0.15), Inches(2.55), Inches(2), Inches(0.5),
                icon, font_size=14, color=color)
    add_textbox(slide, x + Inches(0.15), Inches(2.85), Inches(2), Inches(0.4),
                val, font_size=18, bold=True, color=BLACK)
    add_textbox(slide, x + Inches(0.15), Inches(3.25), Inches(2), Inches(0.35),
                label, font_size=11, color=GRAY)

# ── Rekap Table ──
add_rounded_rect(slide, Inches(2.8), Inches(3.9), Inches(6.5), Inches(3.0),
                 WHITE, LIGHTER)
add_textbox(slide, Inches(3.0), Inches(4.0), Inches(4), Inches(0.4),
            "Rekap Pembayaran per Bulan", font_size=13, bold=True, color=BLACK)
# Table header
headers = ["Periode", "Nominal", "Bayar", "Belum", "Terkumpul"]
for j, h in enumerate(headers):
    add_rect(slide, Inches(3.0 + j*1.2), Inches(4.45), Inches(1.15), Inches(0.35), TEAL)
    add_textbox(slide, Inches(3.05 + j*1.2), Inches(4.45), Inches(1.1), Inches(0.33),
                h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
# Table rows
rows_data = [
    ["Agustus 2026", "Rp 50.000", "38", "12", "Rp 1.900.000"],
    ["Juli 2026", "Rp 50.000", "42", "8", "Rp 2.100.000"],
    ["Juni 2026", "Rp 50.000", "40", "10", "Rp 2.000.000"],
]
for r, row in enumerate(rows_data):
    for j, val in enumerate(row):
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        add_rect(slide, Inches(3.0 + j*1.2), Inches(4.85 + r*0.4), Inches(1.15), Inches(0.38), bg)
        add_textbox(slide, Inches(3.05 + j*1.2), Inches(4.85 + r*0.4), Inches(1.1), Inches(0.36),
                    val, font_size=10, color=BLACK, alignment=PP_ALIGN.CENTER)

# ── Pengeluaran box ──
add_rounded_rect(slide, Inches(9.6), Inches(3.9), Inches(3.2), Inches(3.0),
                 WHITE, LIGHTER)
add_textbox(slide, Inches(9.8), Inches(4.0), Inches(3), Inches(0.4),
            "Pengeluaran Bulan Ini", font_size=13, bold=True, color=BLACK)
peng_items = [
    ("💵 Kas", "Rp 1.000.000", BLUE),
    ("🤝 Sosial", "Rp 1.750.000", GREEN),
    ("🍽️ Konsumsi", "Rp 500.000", ORANGE),
]
for i, (name, amt, clr) in enumerate(peng_items):
    y = Inches(4.5 + i*0.7)
    add_rect(slide, Inches(9.8), y, Inches(2.8), Inches(0.6), LIGHT_GRAY)
    add_textbox(slide, Inches(9.9), y, Inches(1.5), Inches(0.55),
                name, font_size=12, color=clr, bold=True)
    add_textbox(slide, Inches(11.3), y, Inches(1.3), Inches(0.55),
                amt, font_size=12, color=BLACK, alignment=PP_ALIGN.RIGHT)
add_page_number(slide, 6, TOTAL_SLIDES)

# ================================================================
# SLIDE 7 — UI/UX: Monitoring Pembayaran (WIREFRAME)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "UI/UX — Monitoring Pembayaran", "Status iuran warga per periode")

# ── Sidebar (same) ──
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(5.6), SIDEBAR_BG)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.8), DARK_TEAL)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.7),
            "🏠 Iuran RT", font_size=14, bold=True, color=WHITE)
for i, (item, active) in enumerate(sidebar_items):
    y = Inches(2.4 + i * 0.55)
    bg = TEAL if item.startswith("💰") else SIDEBAR_BG
    add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.48), bg)
    add_textbox(slide, Inches(0.7), y, Inches(2), Inches(0.45),
                item, font_size=11, bold=item.startswith("💰"), color=WHITE)

# ── Navbar ──
add_rect(slide, Inches(2.8), Inches(1.5), Inches(10), Inches(0.7), WHITE, LIGHTER, Pt(1))
add_textbox(slide, Inches(3.0), Inches(1.55), Inches(5), Inches(0.5),
            "Monitoring Pembayaran Iuran", font_size=16, bold=True, color=BLACK)
add_rounded_rect(slide, Inches(11.5), Inches(1.6), Inches(1.2), Inches(0.4), BLUE)
add_textbox(slide, Inches(11.5), Inches(1.6), Inches(1.2), Inches(0.4),
            "Export", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ── Card ──
add_rounded_rect(slide, Inches(2.8), Inches(2.4), Inches(10), Inches(4.6), WHITE, LIGHTER)
add_textbox(slide, Inches(3.0), Inches(2.5), Inches(6), Inches(0.4),
            "Status Pembayaran Iuran Warga — Periode Agustus 2026",
            font_size=14, bold=True, color=BLACK)

# Table
pemb_headers = ["#", "Nama Warga", "No. Rumah", "Status", "Tanggal Bayar", "Keterangan"]
col_widths = [0.5, 2.5, 1.2, 1.2, 1.5, 2.8]
for j, h in enumerate(pemb_headers):
    x = Inches(3.0 + sum(col_widths[:j]))
    add_rect(slide, x, Inches(3.0), Inches(col_widths[j]), Inches(0.4), TEAL)
    add_textbox(slide, x + Inches(0.05), Inches(3.0), Inches(col_widths[j]-0.1), Inches(0.38),
                h, font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

pemb_rows = [
    ["1", "Ahmad Fauzi", "01", "Lunas", "2026-08-03", ""],
    ["2", "Budi Santoso", "02", "Belum", "-", "Sudah 2 bln menunggak"],
    ["3", "Citra Lestari", "03", "Lunas", "2026-08-05", ""],
    ["4", "Dedi Kurniawan", "04", "Lunas", "2026-08-07", ""],
    ["5", "Eka Wijaya", "05", "Belum", "-", ""],
]
for r, row in enumerate(pemb_rows):
    y = Inches(3.45 + r * 0.55)
    for j, val in enumerate(row):
        x = Inches(3.0 + sum(col_widths[:j]))
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        add_rect(slide, x, y, Inches(col_widths[j]), Inches(0.5), bg)
        # Status badge coloring
        clr = BLACK
        if j == 3 and val == "Lunas":
            clr = GREEN
        elif j == 3 and val == "Belum":
            clr = RED
        add_textbox(slide, x + Inches(0.05), y, Inches(col_widths[j]-0.1), Inches(0.48),
                    val, font_size=11, color=clr, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 7, TOTAL_SLIDES)

# ================================================================
# SLIDE 8 — UI/UX: Laporan Bulanan (WIREFRAME)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "UI/UX — Laporan Kas Bulanan", "Pemasukan, pengeluaran, saldo akhir")

# Sidebar
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(5.6), SIDEBAR_BG)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.8), DARK_TEAL)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.7),
            "🏠 Iuran RT", font_size=14, bold=True, color=WHITE)
for i, (item, active) in enumerate(sidebar_items):
    y = Inches(2.4 + i * 0.55)
    bg = TEAL if item.startswith("📈") else SIDEBAR_BG
    add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.48), bg)
    add_textbox(slide, Inches(0.7), y, Inches(2), Inches(0.45),
                item, font_size=11, bold=item.startswith("📈"), color=WHITE)

# Navbar
add_rect(slide, Inches(2.8), Inches(1.5), Inches(10), Inches(0.7), WHITE, LIGHTER, Pt(1))
add_textbox(slide, Inches(3.0), Inches(1.55), Inches(5), Inches(0.5),
            "Laporan Kas RT — Agustus 2026", font_size=16, bold=True, color=BLACK)
add_rounded_rect(slide, Inches(11.5), Inches(1.6), Inches(1.2), Inches(0.4), GREEN)
add_textbox(slide, Inches(11.5), Inches(1.6), Inches(1.2), Inches(0.4),
            "Cetak", font_size=11, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Laporan card
add_rounded_rect(slide, Inches(3.5), Inches(2.5), Inches(8.5), Inches(4.5), WHITE, LIGHTER)

# Pemasukan section
add_rect(slide, Inches(3.5), Inches(2.5), Inches(8.5), Inches(0.5), TEAL)
add_textbox(slide, Inches(3.7), Inches(2.52), Inches(3), Inches(0.45),
            "Pemasukan", font_size=14, bold=True, color=WHITE)

add_rect(slide, Inches(3.5), Inches(3.05), Inches(8.5), Inches(0.5), LIGHT_GRAY)
add_textbox(slide, Inches(3.7), Inches(3.07), Inches(5), Inches(0.45),
            "Iuran warga (38 × Rp 50.000)", font_size=13, color=BLACK)
add_textbox(slide, Inches(9.5), Inches(3.07), Inches(2.3), Inches(0.45),
            "Rp 1.900.000", font_size=13, color=BLACK, alignment=PP_ALIGN.RIGHT)

# Pengeluaran header
add_rect(slide, Inches(3.5), Inches(3.6), Inches(8.5), Inches(0.5), TEAL)
add_textbox(slide, Inches(3.7), Inches(3.62), Inches(3), Inches(0.45),
            "Pengeluaran", font_size=14, bold=True, color=WHITE)

peng_laporan = [
    ("💵 Kas — perbaikan pos kamling", "Rp 1.000.000"),
    ("🤝 Sosial — santunan warga sakit", "Rp 1.750.000"),
    ("🍽️ Konsumsi — rapat RT", "Rp 500.000"),
]
for i, (item, amt) in enumerate(peng_laporan):
    y = Inches(4.15 + i * 0.5)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_rect(slide, Inches(3.5), y, Inches(8.5), Inches(0.48), bg)
    add_textbox(slide, Inches(3.7), y, Inches(5), Inches(0.45),
                item, font_size=13, color=BLACK)
    add_textbox(slide, Inches(9.5), y, Inches(2.3), Inches(0.45),
                amt, font_size=13, color=RED, alignment=PP_ALIGN.RIGHT)

# Saldo
add_rect(slide, Inches(3.5), Inches(5.7), Inches(8.5), Inches(0.55), RGBColor(0xEC, 0xFD, 0xF5))
add_textbox(slide, Inches(3.7), Inches(5.72), Inches(5), Inches(0.5),
            "Saldo Akhir", font_size=16, bold=True, color=DARK_TEAL)
add_textbox(slide, Inches(9.0), Inches(5.72), Inches(2.8), Inches(0.5),
            "Rp 10.850.000", font_size=16, bold=True, color=DARK_TEAL, alignment=PP_ALIGN.RIGHT)
add_page_number(slide, 8, TOTAL_SLIDES)

# ================================================================
# SLIDE 9 — UI/UX: Tagihan Warga (WIREFRAME)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "UI/UX — Tagihan Iuran Warga", "Warga melihat tagihan & informasi iuran")

# Sidebar — Warga mode
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(5.6), SIDEBAR_BG)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.8), DARK_TEAL)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.7),
            "🏠 Iuran RT", font_size=14, bold=True, color=WHITE)

warga_sidebar = [
    ("📊 Dashboard", False),
    ("📄 Tagihan Saya", True),
    ("📜 Riwayat Pembayaran", False),
]
for i, (item, active) in enumerate(warga_sidebar):
    y = Inches(2.4 + i * 0.55)
    bg = TEAL if active else SIDEBAR_BG
    add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.48), bg)
    add_textbox(slide, Inches(0.7), y, Inches(2), Inches(0.45),
                item, font_size=11, bold=active, color=WHITE)

# Navbar
add_rect(slide, Inches(2.8), Inches(1.5), Inches(10), Inches(0.7), WHITE, LIGHTER, Pt(1))
add_textbox(slide, Inches(3.0), Inches(1.55), Inches(5), Inches(0.5),
            "Tagihan Iuran", font_size=16, bold=True, color=BLACK)
add_textbox(slide, Inches(11), Inches(1.55), Inches(1.5), Inches(0.5),
            "👤 Warga", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

# ── Tagihan Card (left) ──
add_rounded_rect(slide, Inches(3.0), Inches(2.5), Inches(3.8), Inches(4.2), WHITE, BLUE)
add_rect(slide, Inches(3.0), Inches(2.5), Inches(3.8), Inches(0.08), BLUE)
add_textbox(slide, Inches(3.2), Inches(2.7), Inches(3.4), Inches(0.4),
            "Tagihan Iuran Anda", font_size=16, bold=True, color=BLACK)
# Big amount
add_textbox(slide, Inches(3.2), Inches(3.3), Inches(3.4), Inches(0.8),
            "Rp 50.000", font_size=36, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.2), Inches(4.1), Inches(3.4), Inches(0.4),
            "Periode Agustus 2026", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)
# Status badge
add_rounded_rect(slide, Inches(4.1), Inches(4.6), Inches(1.6), Inches(0.4), RED)
add_textbox(slide, Inches(4.1), Inches(4.6), Inches(1.6), Inches(0.4),
            "Belum dibayar", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
# Button
add_rounded_rect(slide, Inches(3.2), Inches(5.3), Inches(3.4), Inches(0.6), BLUE)
add_textbox(slide, Inches(3.2), Inches(5.3), Inches(3.4), Inches(0.6),
            "💳 Bayar Sekarang", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ── Info Panel (right) ──
add_rounded_rect(slide, Inches(7.2), Inches(2.5), Inches(5.4), Inches(4.2), WHITE, YELLOW_BG)
add_rect(slide, Inches(7.2), Inches(2.5), Inches(5.4), Inches(0.5), YELLOW_BG)
add_textbox(slide, Inches(7.4), Inches(2.52), Inches(4), Inches(0.45),
            "ℹ️ Informasi", font_size=14, bold=True, color=BLACK)
info_text = (
    "Iuran bulan berjalan sebesar Rp 50.000 dengan rincian:\n\n"
    "  💰 Kas       : Rp 25.000\n"
    "  🤝 Sosial    : Rp 15.000\n"
    "  🍽️ Konsumsi : Rp 10.000\n\n"
    "Pembayaran paling lambat tanggal 10 setiap bulan\n"
    "kepada pengurus RT."
)
add_textbox(slide, Inches(7.4), Inches(3.2), Inches(5), Inches(3.2),
            info_text, font_size=13, color=BLACK)
add_page_number(slide, 9, TOTAL_SLIDES)

# ================================================================
# SLIDE 10 — UI/UX: Riwayat Pembayaran Warga (WIREFRAME)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "UI/UX — Riwayat Pembayaran", "Warga melihat histori pembayaran iuran")

# Sidebar
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(5.6), SIDEBAR_BG)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(2.2), Inches(0.8), DARK_TEAL)
add_textbox(slide, Inches(0.7), Inches(1.55), Inches(2), Inches(0.7),
            "🏠 Iuran RT", font_size=14, bold=True, color=WHITE)
for i, (item, active) in enumerate(warga_sidebar):
    y = Inches(2.4 + i * 0.55)
    bg = TEAL if item.startswith("📜") else SIDEBAR_BG
    add_rect(slide, Inches(0.5), y, Inches(2.2), Inches(0.48), bg)
    add_textbox(slide, Inches(0.7), y, Inches(2), Inches(0.45),
                item, font_size=11, bold=item.startswith("📜"), color=WHITE)

# Navbar
add_rect(slide, Inches(2.8), Inches(1.5), Inches(10), Inches(0.7), WHITE, LIGHTER, Pt(1))
add_textbox(slide, Inches(3.0), Inches(1.55), Inches(5), Inches(0.5),
            "Riwayat Pembayaran Iuran Anda", font_size=16, bold=True, color=BLACK)
add_textbox(slide, Inches(11), Inches(1.55), Inches(1.5), Inches(0.5),
            "👤 Warga", font_size=12, color=GRAY, alignment=PP_ALIGN.RIGHT)

# Card
add_rounded_rect(slide, Inches(2.8), Inches(2.4), Inches(10), Inches(4.5), WHITE, LIGHTER)

# Table
hist_headers = ["Periode", "Nominal", "Tanggal Bayar", "Status", "Bukti"]
hist_widths = [2.0, 1.5, 2.0, 1.5, 1.0]
for j, h in enumerate(hist_headers):
    x = Inches(3.2 + sum(hist_widths[:j]))
    add_rect(slide, x, Inches(2.8), Inches(hist_widths[j]), Inches(0.45), TEAL)
    add_textbox(slide, x + Inches(0.05), Inches(2.82), Inches(hist_widths[j]-0.1), Inches(0.4),
                h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

hist_rows = [
    ["Juli 2026", "Rp 50.000", "2026-07-04", "Lunas", "👁"],
    ["Juni 2026", "Rp 50.000", "2026-06-05", "Lunas", "👁"],
    ["Mei 2026", "Rp 50.000", "2026-05-06", "Lunas", "👁"],
    ["April 2026", "Rp 50.000", "-", "Belum Bayar", "-"],
]
for r, row in enumerate(hist_rows):
    y = Inches(3.3 + r * 0.6)
    for j, val in enumerate(row):
        x = Inches(3.2 + sum(hist_widths[:j]))
        bg = LIGHT_GRAY if r % 2 == 0 else WHITE
        add_rect(slide, x, y, Inches(hist_widths[j]), Inches(0.55), bg)
        clr = BLACK
        if j == 3 and val == "Lunas":
            clr = GREEN
        elif j == 3 and val == "Belum Bayar":
            clr = RED
        add_textbox(slide, x + Inches(0.05), y, Inches(hist_widths[j]-0.1), Inches(0.53),
                    val, font_size=12, color=clr, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 10, TOTAL_SLIDES)

# ================================================================
# SLIDE 11 — Task 3: Design Database (ERD)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Task 3 — Design Database", "6 Tabel Relasional • MySQL/MariaDB")

# Tabel info
tbl_data = [
    ["Tabel", "Fungsi"],
    ["users", "Akun login (pengurus & warga)"],
    ["warga", "Data warga & status aktif"],
    ["pengaturan_iuran", "Nominal iuran berjalan (historis)"],
    ["pembayaran", "Pembayaran warga per periode"],
    ["kategori_pengeluaran", "Kas, Sosial, Konsumsi"],
    ["pengeluaran", "Transaksi pengeluaran RT"],
]
tbl_shape = slide.shapes.add_table(7, 2, Inches(0.8), Inches(1.7), Inches(5.5), Inches(4)).table
tbl_shape.columns[0].width = Inches(2.5)
tbl_shape.columns[1].width = Inches(3)
for r, row in enumerate(tbl_data):
    for c, val in enumerate(row):
        cell = tbl_shape.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = 'Segoe UI'
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL if r == 0 else (LIGHT_GRAY if r % 2 == 1 else WHITE)

# ERD visual
add_rounded_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(3.5),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(7.0), Inches(1.8), Inches(4), Inches(0.4),
            "Entity Relationship Diagram", font_size=16, bold=True, color=TEAL)

erd_lines = [
    ("users", 7.2, 2.4, BLUE),
    ("warga", 9.2, 2.4, TEAL),
    ("pembayaran", 11.2, 2.4, GREEN),
    ("pengaturan_iuran", 8.0, 3.5, ORANGE),
    ("kategori_pengeluaran", 10.5, 3.5, RED),
    ("pengeluaran", 10.5, 4.3, RED),
]
for name, x, y, clr in erd_lines:
    add_rounded_rect(slide, Inches(x), Inches(y), Inches(1.6), Inches(0.5), clr)
    add_textbox(slide, Inches(x), Inches(y), Inches(1.6), Inches(0.5),
                name, font_size=9, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Relationship notes
add_textbox(slide, Inches(7.0), Inches(5.0), Inches(5.4), Inches(2),
            "• warga 1→* pembayaran, UNIQUE(warga_id, periode)\n"
            "• pembayaran.iuran_id → nominal historis\n"
            "• kategori_pengeluaran 1→* pengeluaran\n"
            "• Perubahan iuran = record baru (tidak ubah lama)",
            font_size=12, color=BLACK)
add_page_number(slide, 11, TOTAL_SLIDES)

# ================================================================
# SLIDE 12 — Rencana Pengembangan
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_slide_title(slide, "Demo & Rencana Pengembangan")

# Demo section
add_rounded_rect(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.0),
                 LIGHT_GRAY, LIGHTER)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4), Inches(0.5),
            "Demo Aplikasi", font_size=20, bold=True, color=TEAL)
demo_items = [
    "• Menjalankan: php spark serve → localhost:8080",
    "• Dashboard pengurus: info kas, rekap bulanan",
    "• Monitoring pembayaran warga",
    "• Laporan kas bulanan",
    "• Tagihan iuran (peran warga)",
    "• Riwayat pembayaran (peran warga)",
]
for i, item in enumerate(demo_items):
    add_textbox(slide, Inches(1.2), Inches(2.7 + i*0.55), Inches(5.2), Inches(0.5),
                item, font_size=14, color=BLACK)

# Rencana pengembangan
add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.0),
                 RGBColor(0xEC, 0xFD, 0xF5), GREEN)
add_textbox(slide, Inches(7.4), Inches(1.9), Inches(4), Inches(0.5),
            "Rencana Pengembangan", font_size=20, bold=True, color=DARK_TEAL)
dev_items = [
    "🔐  Autentikasi & otorisasi (login per peran)",
    "🗃️  CRUD warga, pembayaran, pengeluaran",
    "     (terhubung database)",
    "📄  Cetak laporan PDF",
    "🔔  Notifikasi tagihan otomatis",
    "📱  Responsive mobile view",
    "📊  Chart/grafik pembayaran",
]
for i, item in enumerate(dev_items):
    add_textbox(slide, Inches(7.4), Inches(2.7 + i*0.55), Inches(5), Inches(0.5),
                item, font_size=14, color=BLACK)
add_page_number(slide, 12, TOTAL_SLIDES)

# ================================================================
# SLIDE 13 — Terima Kasih
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_TEAL)
add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
            "Terima Kasih", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.7),
            "Pertanyaan & saran dipersilakan", font_size=22, color=RGBColor(0xA7, 0xF3, 0xD0),
            alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.04), WHITE)
add_page_number(slide, 13, TOTAL_SLIDES)

# ── Save ──
output_path = r"C:\Users\zamih\PWL\Aplikasi_RT\presentasi\Presentasi_Aplikasi_RT_UIUX.pptx"
prs.save(output_path)
print(f"✅ Presentasi berhasil disimpan ke:\n{output_path}")
